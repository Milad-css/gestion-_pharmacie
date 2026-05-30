from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Color palette ─────────────────────────────────────────────
C_GREEN_DARK  = RGBColor(0x1A, 0x6B, 0x3C)   # deep pharmacy green
C_GREEN_MID   = RGBColor(0x27, 0xAE, 0x60)   # medium green
C_GREEN_LIGHT = RGBColor(0xD5, 0xF5, 0xE3)   # very light green bg
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK        = RGBColor(0x17, 0x20, 0x2A)   # near-black text
C_GRAY        = RGBColor(0x5D, 0x6D, 0x7E)   # secondary text
C_ACCENT      = RGBColor(0x21, 0x8C, 0x74)   # teal accent
C_LIGHT_BG    = RGBColor(0xF4, 0xF6, 0xF7)   # slide background


# ── Helpers ───────────────────────────────────────────────────
def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=C_DARK,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_bullet_box(slide, items, left, top, width, height,
                   title=None, title_color=C_GREEN_DARK, item_size=14,
                   bullet_color=C_GREEN_MID, bg=None, border=None):
    if bg:
        add_rect(slide, left, top, width, height, bg, border)
    if title:
        add_textbox(slide, title, left + Inches(0.15), top + Inches(0.1),
                    width - Inches(0.3), Inches(0.4),
                    font_size=15, bold=True, color=title_color)
        t_off = Inches(0.45)
    else:
        t_off = Inches(0.1)

    txb = slide.shapes.add_textbox(
        left + Inches(0.15), top + t_off,
        width - Inches(0.3), height - t_off - Inches(0.1)
    )
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(item_size)
        p.font.color.rgb = C_DARK
    return txb


# ── Build presentation ────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]   # completely blank


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, C_GREEN_DARK)

# big diagonal accent
acc = add_rect(sl, Inches(7.8), Inches(-0.5), Inches(6), Inches(9),
               C_GREEN_MID)

# white card
add_rect(sl, Inches(0.6), Inches(1.5), Inches(7.5), Inches(4.5), C_WHITE)

# cross icon (pharmacy symbol) - just styled text
add_textbox(sl, "✚", Inches(0.9), Inches(1.6), Inches(1.2), Inches(1.2),
            font_size=52, bold=True, color=C_GREEN_MID, align=PP_ALIGN.CENTER)

add_textbox(sl, "Gestion Pharmacie",
            Inches(1.9), Inches(1.65), Inches(6), Inches(1.0),
            font_size=38, bold=True, color=C_GREEN_DARK)

add_textbox(sl, "Système de gestion d'une pharmacie en ligne",
            Inches(1.9), Inches(2.6), Inches(6), Inches(0.7),
            font_size=17, color=C_GRAY, italic=True)

add_rect(sl, Inches(0.6), Inches(3.4), Inches(7.5), Inches(0.04), C_GREEN_MID)

add_textbox(sl, "Laravel 13  •  React  •  Vite  •  Tailwind CSS  •  SQLite",
            Inches(0.7), Inches(3.55), Inches(7.3), Inches(0.5),
            font_size=14, color=C_ACCENT, align=PP_ALIGN.CENTER)

add_textbox(sl, "API REST  |  Authentification Sanctum  |  Rôles Admin / Client",
            Inches(0.7), Inches(4.05), Inches(7.3), Inches(0.5),
            font_size=13, color=C_GRAY, align=PP_ALIGN.CENTER)

add_textbox(sl, "2025", Inches(0.7), Inches(5.5), Inches(7.3), Inches(0.5),
            font_size=13, color=C_WHITE, align=PP_ALIGN.CENTER)

# right side decorative text
add_textbox(sl, "💊", Inches(9.5), Inches(2.5), Inches(2), Inches(2),
            font_size=80, align=PP_ALIGN.CENTER, color=C_WHITE)


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — SOMMAIRE
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, C_LIGHT_BG)
add_rect(sl, 0, 0, Inches(13.33), Inches(1.1), C_GREEN_DARK)
add_textbox(sl, "Sommaire", Inches(0.4), Inches(0.15), Inches(12), Inches(0.8),
            font_size=30, bold=True, color=C_WHITE)

sections = [
    ("01", "Vue d'ensemble du projet",      "Objectifs et contexte"),
    ("02", "Stack technique",               "Technologies utilisées"),
    ("03", "Architecture du système",       "Backend + Frontend"),
    ("04", "Base de données",               "Modèles & relations"),
    ("05", "API REST — Backend",            "Routes & contrôleurs"),
    ("06", "Interface Client (Frontend)",   "Pages & fonctionnalités"),
    ("07", "Interface Admin",               "Tableau de bord & gestion"),
    ("08", "Sécurité & Authentification",   "Sanctum, middleware, rôles"),
    ("09", "Lancer le projet",              "Installation & démarrage"),
]

cols = 3
rows = 3
box_w = Inches(3.9)
box_h = Inches(1.5)
gap_x = Inches(0.3)
gap_y = Inches(0.3)
start_x = Inches(0.35)
start_y = Inches(1.3)

for i, (num, title, sub) in enumerate(sections):
    col = i % cols
    row = i // cols
    lx = start_x + col * (box_w + gap_x)
    ty = start_y + row * (box_h + gap_y)
    add_rect(sl, lx, ty, box_w, box_h, C_WHITE, C_GREEN_MID)
    add_rect(sl, lx, ty, Inches(0.55), box_h, C_GREEN_MID)
    add_textbox(sl, num, lx + Inches(0.05), ty + Inches(0.4),
                Inches(0.55), Inches(0.6), font_size=16, bold=True,
                color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, title, lx + Inches(0.65), ty + Inches(0.15),
                box_w - Inches(0.75), Inches(0.6),
                font_size=14, bold=True, color=C_GREEN_DARK)
    add_textbox(sl, sub, lx + Inches(0.65), ty + Inches(0.75),
                box_w - Inches(0.75), Inches(0.6),
                font_size=12, color=C_GRAY, italic=True)


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — VUE D'ENSEMBLE
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, C_LIGHT_BG)
add_rect(sl, 0, 0, Inches(13.33), Inches(1.1), C_GREEN_DARK)
add_textbox(sl, "01 — Vue d'ensemble du projet",
            Inches(0.4), Inches(0.15), Inches(12), Inches(0.8),
            font_size=28, bold=True, color=C_WHITE)

add_textbox(sl, "Contexte",
            Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.5),
            font_size=18, bold=True, color=C_GREEN_DARK)
add_textbox(sl,
    "Application web complète pour la gestion d'une pharmacie en ligne. "
    "Elle permet à une pharmacie de présenter son catalogue, "
    "gérer les commandes clients, et administrer ses produits et stocks depuis un tableau de bord dédié.",
    Inches(0.4), Inches(1.65), Inches(12.5), Inches(0.9),
    font_size=14, color=C_DARK)

# Two goal columns
goals_client = [
    "Consulter le catalogue de médicaments",
    "Rechercher et filtrer par catégorie",
    "Ajouter des produits au panier",
    "Passer et suivre ses commandes",
    "Gérer son profil personnel",
]
goals_admin = [
    "Tableau de bord avec statistiques",
    "CRUD complet sur produits et catégories",
    "Upload d'images pour les produits",
    "Gestion et suivi de toutes les commandes",
    "Changement de statut des commandes",
]

add_bullet_box(sl, goals_client, Inches(0.4), Inches(2.7),
               Inches(5.9), Inches(3.5),
               title="Fonctionnalités Client", bg=C_WHITE, border=C_GREEN_LIGHT)
add_bullet_box(sl, goals_admin, Inches(6.9), Inches(2.7),
               Inches(6.0), Inches(3.5),
               title="Fonctionnalités Admin", bg=C_WHITE, border=C_GREEN_LIGHT)

add_rect(sl, Inches(6.4), Inches(2.7), Inches(0.06), Inches(3.5), C_GREEN_MID)


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — STACK TECHNIQUE
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, C_LIGHT_BG)
add_rect(sl, 0, 0, Inches(13.33), Inches(1.1), C_GREEN_DARK)
add_textbox(sl, "02 — Stack technique",
            Inches(0.4), Inches(0.15), Inches(12), Inches(0.8),
            font_size=28, bold=True, color=C_WHITE)

tech_groups = [
    ("Backend",   C_GREEN_DARK, [
        "PHP 8 / Laravel 13",
        "Laravel Sanctum (auth par token)",
        "Eloquent ORM",
        "SQLite (base de données légère)",
        "API REST (JSON)",
        "Middleware Admin personnalisé",
    ]),
    ("Frontend",  C_ACCENT, [
        "React 18",
        "Vite (bundler ultra-rapide)",
        "Tailwind CSS v4",
        "Axios (client HTTP)",
        "React Router v6",
        "Context API (Auth + Cart)",
    ]),
    ("Outils",    C_GRAY, [
        "Composer (dépendances PHP)",
        "npm (dépendances JS)",
        "Git (contrôle de version)",
        "Storage symlink (images)",
        "CORS configuré",
        "Seeders (données initiales)",
    ]),
]

bw = Inches(3.9)
for i, (grp, col, items) in enumerate(tech_groups):
    lx = Inches(0.35) + i * (bw + Inches(0.3))
    add_rect(sl, lx, Inches(1.25), bw, Inches(0.5), col)
    add_textbox(sl, grp, lx, Inches(1.25), bw, Inches(0.5),
                font_size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, lx, Inches(1.75), bw, Inches(4.5), C_WHITE, C_GREEN_LIGHT)
    add_bullet_box(sl, items, lx, Inches(1.75), bw, Inches(4.5),
                   item_size=14, bg=None)


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — ARCHITECTURE
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, C_LIGHT_BG)
add_rect(sl, 0, 0, Inches(13.33), Inches(1.1), C_GREEN_DARK)
add_textbox(sl, "03 — Architecture du système",
            Inches(0.4), Inches(0.15), Inches(12), Inches(0.8),
            font_size=28, bold=True, color=C_WHITE)

# Architecture diagram boxes
boxes = [
    (Inches(0.4),  Inches(1.3), Inches(3.0), Inches(2.0), C_WHITE, "Navigateur\n(Client Web)", C_ACCENT),
    (Inches(5.0),  Inches(1.3), Inches(3.5), Inches(2.0), C_WHITE, "Frontend\nReact + Vite\nlocalhost:5173", C_GREEN_MID),
    (Inches(9.3),  Inches(1.3), Inches(3.6), Inches(2.0), C_WHITE, "Backend\nLaravel API\nlocalhost:8000", C_GREEN_DARK),
    (Inches(5.0),  Inches(4.3), Inches(3.5), Inches(1.8), C_WHITE, "Contextes React\nAuthContext / CartContext", C_ACCENT),
    (Inches(9.3),  Inches(4.3), Inches(3.6), Inches(1.8), C_WHITE, "Base de données\nSQLite\n(6 tables)", C_GREEN_DARK),
]
for (lx, ty, bw, bh, bg, lbl, tc) in boxes:
    add_rect(sl, lx, ty, bw, bh, bg, tc)
    add_rect(sl, lx, ty, bw, Inches(0.08), tc)
    add_textbox(sl, lbl, lx, ty + Inches(0.1), bw, bh - Inches(0.1),
                font_size=14, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

# Arrows (just text labels on lines)
arr_labels = [
    (Inches(3.5),  Inches(2.1), "HTTP / HTTPS →"),
    (Inches(8.6),  Inches(2.1), "Axios / API REST →"),
    (Inches(6.5),  Inches(3.3), "↕  Sync"),
    (Inches(10.8), Inches(3.3), "↕  Eloquent ORM"),
]
for (lx, ty, lbl) in arr_labels:
    add_textbox(sl, lbl, lx, ty, Inches(1.5), Inches(0.5),
                font_size=11, italic=True, color=C_GRAY, align=PP_ALIGN.CENTER)

add_textbox(sl,
    "Architecture découplée : le frontend React communique exclusivement avec le backend via l'API REST. "
    "Le backend gère l'authentification (Sanctum), la logique métier et la persistance des données.",
    Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.9),
    font_size=12, italic=True, color=C_GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 6 — BASE DE DONNÉES
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, C_LIGHT_BG)
add_rect(sl, 0, 0, Inches(13.33), Inches(1.1), C_GREEN_DARK)
add_textbox(sl, "04 — Base de données",
            Inches(0.4), Inches(0.15), Inches(12), Inches(0.8),
            font_size=28, bold=True, color=C_WHITE)

tables = [
    ("users",       ["id, name, email", "password, role", "telephone, adresse"]),
    ("categories",  ["id, nom", "description"]),
    ("products",    ["id, category_id", "nom, description", "prix, stock", "image, actif", "date_expiration"]),
    ("orders",      ["id, user_id", "total, statut", "adresse, telephone", "notes"]),
    ("order_items", ["id, order_id", "product_id", "quantite", "prix_unitaire"]),
    ("cart_items",  ["id, user_id", "product_id", "quantite"]),
]

tw = Inches(2.05)
th = Inches(2.4)
gap = Inches(0.22)
sx = Inches(0.3)
sy = Inches(1.3)

for i, (tname, cols) in enumerate(tables):
    col = i % 3
    row = i // 3
    lx = sx + col * (tw + gap)
    ty = sy + row * (th + gap)
    add_rect(sl, lx, ty, tw, th, C_WHITE, C_GREEN_MID)
    add_rect(sl, lx, ty, tw, Inches(0.42), C_GREEN_MID)
    add_textbox(sl, tname, lx, ty, tw, Inches(0.42),
                font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    for j, col_txt in enumerate(cols):
        add_textbox(sl, col_txt,
                    lx + Inches(0.08), ty + Inches(0.45) + j * Inches(0.38),
                    tw - Inches(0.16), Inches(0.36),
                    font_size=11, color=C_DARK)

# Relations
add_rect(sl, Inches(6.6), Inches(1.3), Inches(6.6), Inches(5.3), C_WHITE, C_GREEN_LIGHT)
add_textbox(sl, "Relations Eloquent", Inches(6.7), Inches(1.35), Inches(6.4), Inches(0.5),
            font_size=15, bold=True, color=C_GREEN_DARK)
rels = [
    "User  →  hasMany Orders",
    "User  →  hasMany CartItems",
    "Category  →  hasMany Products",
    "Product  →  belongsTo Category",
    "Order  →  hasMany OrderItems",
    "OrderItem  →  belongsTo Product",
    "CartItem  →  belongsTo Product",
    "Statuts commande : en_attente,",
    "  confirmee, livree, annulee",
]
add_bullet_box(sl, rels, Inches(6.6), Inches(1.8), Inches(6.6), Inches(4.6),
               item_size=13, bg=None)


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — API REST
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, C_LIGHT_BG)
add_rect(sl, 0, 0, Inches(13.33), Inches(1.1), C_GREEN_DARK)
add_textbox(sl, "05 — API REST — Backend",
            Inches(0.4), Inches(0.15), Inches(12), Inches(0.8),
            font_size=28, bold=True, color=C_WHITE)

route_groups = [
    ("Publiques (sans auth)", C_GRAY, [
        "POST  /api/register",
        "POST  /api/login",
        "GET   /api/categories",
        "GET   /api/categories/{id}",
        "GET   /api/products",
        "GET   /api/products/{id}",
    ]),
    ("Authentifiées (Sanctum)", C_ACCENT, [
        "POST  /api/logout",
        "GET   /api/profile",
        "PUT   /api/profile",
        "GET/POST  /api/cart",
        "PUT/DELETE  /api/cart/{id}",
        "DELETE  /api/cart",
        "GET/POST  /api/orders",
        "GET   /api/orders/{id}",
    ]),
    ("Admin uniquement", C_GREEN_DARK, [
        "POST  /api/categories",
        "PUT/DELETE  /api/categories/{id}",
        "POST  /api/products",
        "POST  /api/products/{id}",
        "DELETE  /api/products/{id}",
        "GET   /api/admin/orders",
        "PUT   /api/admin/orders/{id}",
    ]),
]

bw = Inches(4.0)
for i, (grp, col, routes) in enumerate(route_groups):
    lx = Inches(0.3) + i * (bw + Inches(0.25))
    add_rect(sl, lx, Inches(1.25), bw, Inches(0.45), col)
    add_textbox(sl, grp, lx, Inches(1.25), bw, Inches(0.45),
                font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, lx, Inches(1.7), bw, Inches(4.65), C_WHITE, C_GREEN_LIGHT)
    add_bullet_box(sl, routes, lx, Inches(1.7), bw, Inches(4.65),
                   item_size=12, bg=None)

add_textbox(sl,
    "25 routes au total  •  Réponses JSON  •  Pagination 20 articles/page  •  Upload image (max 2 Mo)",
    Inches(0.3), Inches(6.5), Inches(12.7), Inches(0.5),
    font_size=12, italic=True, color=C_GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 8 — FRONTEND CLIENT
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, C_LIGHT_BG)
add_rect(sl, 0, 0, Inches(13.33), Inches(1.1), C_GREEN_DARK)
add_textbox(sl, "06 — Interface Client (Frontend)",
            Inches(0.4), Inches(0.15), Inches(12), Inches(0.8),
            font_size=28, bold=True, color=C_WHITE)

pages = [
    ("Accueil / Catalogue",  ["Catalogue paginé (20 produits/page)", "Barre de recherche en temps réel", "Filtre par catégorie", "Bouton Ajouter au panier"]),
    ("Détail Produit",       ["Tous les détails du produit", "Sélecteur de quantité (1 → stock)", "Bouton Ajouter au panier"]),
    ("Panier",               ["Liste articles avec modification", "Calcul du total en temps réel", "Formulaire adresse / téléphone", "Validation et envoi de commande"]),
    ("Mes Commandes",        ["Liste de toutes ses commandes", "Accordéon détail par commande", "Badge de statut coloré"]),
    ("Profil",               ["Modifier nom, tel, adresse", "Changer mot de passe (optionnel)", "Feedback succès / erreur"]),
    ("Login / Inscription",  ["Formulaire de connexion", "Inscription complète", "Erreurs champ par champ", "Redirection automatique rôle"]),
]

bw = Inches(3.9)
bh = Inches(2.5)
gap = Inches(0.25)
sx = Inches(0.35)
sy = Inches(1.25)

for i, (title, items) in enumerate(pages):
    col = i % 3
    row = i // 2
    # Place 3 per row, 2 rows
    col_idx = i % 3
    row_idx = i // 3
    lx = sx + col_idx * (bw + gap)
    ty = sy + row_idx * (bh + gap)
    add_rect(sl, lx, ty, bw, bh, C_WHITE, C_GREEN_LIGHT)
    add_rect(sl, lx, ty, bw, Inches(0.4), C_GREEN_MID)
    add_textbox(sl, title, lx, ty, bw, Inches(0.4),
                font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_bullet_box(sl, items, lx, ty + Inches(0.4), bw, bh - Inches(0.4),
                   item_size=12, bg=None)


# ════════════════════════════════════════════════════════════════
# SLIDE 9 — INTERFACE ADMIN
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, C_LIGHT_BG)
add_rect(sl, 0, 0, Inches(13.33), Inches(1.1), C_GREEN_DARK)
add_textbox(sl, "07 — Interface Admin",
            Inches(0.4), Inches(0.15), Inches(12), Inches(0.8),
            font_size=28, bold=True, color=C_WHITE)

admin_pages = [
    ("Dashboard",           C_GREEN_DARK, [
        "Nombre de produits",
        "Nombre de catégories",
        "Total commandes",
        "Commandes en attente",
        "Liens rapides vers les sections",
    ]),
    ("Gestion Produits",    C_ACCENT, [
        "Formulaire CRUD complet",
        "Upload d'image produit",
        "Liste avec aperçu image",
        "Stock et statut actif/inactif",
        "Modification inline (pré-rempli)",
        "Suppression avec image",
    ]),
    ("Gestion Catégories",  C_GREEN_MID, [
        "Formulaire CRUD simple",
        "Compteur de produits / catégorie",
        "Modification et suppression",
    ]),
    ("Gestion Commandes",   C_GREEN_DARK, [
        "Toutes les commandes (paginées)",
        "Accordéon détail de chaque commande",
        "Boutons changement de statut",
        "4 statuts : en_attente, confirmée,",
        "  livrée, annulée",
    ]),
]

bw = Inches(3.0)
bh = Inches(4.5)
gap = Inches(0.26)
sx = Inches(0.35)

for i, (title, col, items) in enumerate(admin_pages):
    lx = sx + i * (bw + gap)
    add_rect(sl, lx, Inches(1.25), bw, bh, C_WHITE, C_GREEN_LIGHT)
    add_rect(sl, lx, Inches(1.25), bw, Inches(0.45), col)
    add_textbox(sl, title, lx, Inches(1.25), bw, Inches(0.45),
                font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_bullet_box(sl, items, lx, Inches(1.7), bw, bh - Inches(0.45),
                   item_size=12.5, bg=None)

add_textbox(sl,
    "Accès réservé aux comptes avec role = 'admin'  •  AdminRoute bloque tout accès non autorisé",
    Inches(0.3), Inches(6.0), Inches(12.7), Inches(0.5),
    font_size=12, italic=True, color=C_GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 10 — SÉCURITÉ
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, C_LIGHT_BG)
add_rect(sl, 0, 0, Inches(13.33), Inches(1.1), C_GREEN_DARK)
add_textbox(sl, "08 — Sécurité & Authentification",
            Inches(0.4), Inches(0.15), Inches(12), Inches(0.8),
            font_size=28, bold=True, color=C_WHITE)

sec_items = [
    ("Laravel Sanctum", [
        "Authentification par token Bearer",
        "Token stocké dans localStorage",
        "Envoyé automatiquement par Axios",
        "Révocation à la déconnexion",
    ]),
    ("Middleware Admin", [
        "Vérifie role === 'admin'",
        "Retourne HTTP 403 sinon",
        "Alias 'admin' dans bootstrap/app.php",
        "Appliqué sur toutes les routes admin",
    ]),
    ("Validation & Sécurité", [
        "Validation Laravel sur toutes les entrées",
        "Hash bcrypt des mots de passe",
        "Vérification d'appartenance (panier)",
        "Transactions DB pour les commandes",
        "Vérification de stock avant commande",
        "CORS : origines autorisées uniquement",
    ]),
    ("Route Guards React", [
        "ProtectedRoute : redirige si non connecté",
        "AdminRoute : redirige si non admin",
        "AuthContext vérifie le token au démarrage",
        "Erreurs 401/403 gérées côté client",
    ]),
]

bw = Inches(3.0)
bh = Inches(4.5)
gap = Inches(0.26)
sx = Inches(0.35)

for i, (title, items) in enumerate(sec_items):
    lx = sx + i * (bw + gap)
    add_rect(sl, lx, Inches(1.25), bw, bh, C_WHITE, C_GREEN_LIGHT)
    add_rect(sl, lx, Inches(1.25), bw, Inches(0.45), C_GREEN_DARK)
    add_textbox(sl, title, lx, Inches(1.25), bw, Inches(0.45),
                font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_bullet_box(sl, items, lx, Inches(1.7), bw, bh - Inches(0.45),
                   item_size=12.5, bg=None)


# ════════════════════════════════════════════════════════════════
# SLIDE 11 — LANCER LE PROJET
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, C_LIGHT_BG)
add_rect(sl, 0, 0, Inches(13.33), Inches(1.1), C_GREEN_DARK)
add_textbox(sl, "09 — Lancer le projet",
            Inches(0.4), Inches(0.15), Inches(12), Inches(0.8),
            font_size=28, bold=True, color=C_WHITE)

# Backend box
add_rect(sl, Inches(0.4), Inches(1.25), Inches(5.9), Inches(4.5), C_WHITE, C_GREEN_MID)
add_rect(sl, Inches(0.4), Inches(1.25), Inches(5.9), Inches(0.5), C_GREEN_DARK)
add_textbox(sl, "Backend — Laravel (port 8000)",
            Inches(0.4), Inches(1.25), Inches(5.9), Inches(0.5),
            font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

backend_steps = [
    "cd backend",
    "composer install",
    "cp .env.example .env",
    "php artisan key:generate",
    "php artisan migrate --seed",
    "php artisan storage:link",
    "php artisan serve",
]
for j, step in enumerate(backend_steps):
    bg = C_DARK if j % 2 == 0 else RGBColor(0x2C, 0x3E, 0x50)
    add_rect(sl, Inches(0.4), Inches(1.75) + j * Inches(0.38), Inches(5.9), Inches(0.37), bg)
    add_textbox(sl, f"$ {step}",
                Inches(0.55), Inches(1.77) + j * Inches(0.38), Inches(5.6), Inches(0.34),
                font_size=12, color=C_GREEN_MID, bold=False)

# Frontend box
add_rect(sl, Inches(7.0), Inches(1.25), Inches(5.9), Inches(4.5), C_WHITE, C_ACCENT)
add_rect(sl, Inches(7.0), Inches(1.25), Inches(5.9), Inches(0.5), C_ACCENT)
add_textbox(sl, "Frontend — React / Vite (port 5173)",
            Inches(7.0), Inches(1.25), Inches(5.9), Inches(0.5),
            font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

frontend_steps = [
    "cd frontend",
    "npm install",
    "npm run dev",
]
for j, step in enumerate(frontend_steps):
    bg = C_DARK if j % 2 == 0 else RGBColor(0x2C, 0x3E, 0x50)
    add_rect(sl, Inches(7.0), Inches(1.75) + j * Inches(0.38), Inches(5.9), Inches(0.37), bg)
    add_textbox(sl, f"$ {step}",
                Inches(7.15), Inches(1.77) + j * Inches(0.38), Inches(5.6), Inches(0.34),
                font_size=12, color=C_GREEN_MID, bold=False)

# Accounts
add_rect(sl, Inches(0.4), Inches(5.9), Inches(12.5), Inches(1.4), C_WHITE, C_GREEN_LIGHT)
add_textbox(sl, "Comptes de test (après seed)",
            Inches(0.5), Inches(5.95), Inches(4), Inches(0.4),
            font_size=13, bold=True, color=C_GREEN_DARK)
add_textbox(sl, "Admin :  admin@pharmacie.com  /  password",
            Inches(0.5), Inches(6.35), Inches(5.5), Inches(0.4),
            font_size=12, color=C_DARK)
add_textbox(sl, "Client :  client@test.com  /  password",
            Inches(6.5), Inches(6.35), Inches(5.5), Inches(0.4),
            font_size=12, color=C_DARK)
add_textbox(sl, "Accéder à :  http://localhost:5173",
            Inches(0.5), Inches(6.75), Inches(12), Inches(0.4),
            font_size=12, italic=True, color=C_ACCENT)


# ════════════════════════════════════════════════════════════════
# SLIDE 12 — CONCLUSION
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, C_GREEN_DARK)

# decorative accent
add_rect(sl, Inches(8.5), Inches(-0.5), Inches(5.5), Inches(9), C_GREEN_MID)

# white card
add_rect(sl, Inches(0.5), Inches(1.3), Inches(8.0), Inches(5.5), C_WHITE)

add_textbox(sl, "✚", Inches(0.7), Inches(1.45), Inches(1.0), Inches(0.9),
            font_size=42, bold=True, color=C_GREEN_MID, align=PP_ALIGN.CENTER)

add_textbox(sl, "Bilan du projet",
            Inches(1.6), Inches(1.4), Inches(6.5), Inches(0.8),
            font_size=30, bold=True, color=C_GREEN_DARK)

add_rect(sl, Inches(0.5), Inches(2.25), Inches(8.0), Inches(0.05), C_GREEN_MID)

achievements = [
    "Application full-stack complète et fonctionnelle",
    "Architecture découplée Backend / Frontend",
    "Authentification sécurisée avec Laravel Sanctum",
    "Deux rôles distincts : Admin et Client",
    "25 endpoints API REST bien structurés",
    "CRUD complet avec upload d'images",
    "Gestion des stocks et transactions DB",
    "Interface responsive avec Tailwind CSS",
]
add_bullet_box(sl, achievements, Inches(0.5), Inches(2.35),
               Inches(8.0), Inches(4.3), item_size=14, bg=None,
               bullet_color=C_GREEN_MID)

# Right side
add_textbox(sl, "Technologies\nmaîtrisées",
            Inches(9.0), Inches(2.0), Inches(3.8), Inches(1.0),
            font_size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

techs = ["Laravel 13", "React 18", "Tailwind CSS v4",
         "Sanctum Auth", "Eloquent ORM", "Vite", "SQLite", "REST API"]
for i, t in enumerate(techs):
    col_i = i % 2
    row_i = i // 2
    lx = Inches(9.0) + col_i * Inches(1.95)
    ty = Inches(3.1) + row_i * Inches(0.75)
    add_rect(sl, lx, ty, Inches(1.8), Inches(0.6), C_WHITE)
    add_textbox(sl, t, lx, ty, Inches(1.8), Inches(0.6),
                font_size=13, bold=True, color=C_GREEN_DARK, align=PP_ALIGN.CENTER)

add_textbox(sl, "Merci !",
            Inches(9.5), Inches(6.5), Inches(3.0), Inches(0.7),
            font_size=24, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────
out = r"d:\gestion _pharmacie\Presentation_Gestion_Pharmacie.pptx"
prs.save(out)
print(f"OK - Fichier cree : {out}")
